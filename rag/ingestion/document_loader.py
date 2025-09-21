"""Document loader for various file formats"""

import os
import mimetypes
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
import json
import logging

# Document processing imports
import pypdf
from pypdf import PdfReader
import docx
from docx import Document
import pptx
from pptx import Presentation
import openpyxl
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import markdown
import pandas as pd

from rich.console import Console

console = Console()
logger = logging.getLogger(__name__)

class DocumentLoader:
    """Load and extract text from various document formats"""

    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.doc': self._load_docx,
            '.pptx': self._load_pptx,
            '.ppt': self._load_pptx,
            '.xlsx': self._load_excel,
            '.xls': self._load_excel,
            '.csv': self._load_csv,
            '.txt': self._load_text,
            '.md': self._load_markdown,
            '.html': self._load_html,
            '.htm': self._load_html,
            '.json': self._load_json,
            '.py': self._load_text,
            '.js': self._load_text,
            '.cpp': self._load_text,
            '.java': self._load_text,
            '.go': self._load_text,
            '.rs': self._load_text,
            '.sql': self._load_text,
            '.yaml': self._load_text,
            '.yml': self._load_text,
            '.xml': self._load_html,
        }

    def load_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """Load a document and extract its content"""
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = file_path.suffix.lower()

        if extension not in self.supported_extensions:
            console.print(f"[yellow]Unsupported file type: {extension}. Trying as text...[/yellow]")
            loader_func = self._load_text
        else:
            loader_func = self.supported_extensions[extension]

        try:
            content = loader_func(file_path)

            # Get file metadata
            stat = file_path.stat()
            metadata = {
                'filename': file_path.name,
                'filepath': str(file_path),
                'extension': extension,
                'size_bytes': stat.st_size,
                'modified_time': stat.st_mtime,
                'file_type': extension[1:] if extension else 'unknown'
            }

            return {
                'content': content,
                'metadata': metadata,
                'source': str(file_path)
            }

        except Exception as e:
            logger.error(f"Error loading {file_path}: {e}")
            return {
                'content': "",
                'metadata': {'filename': file_path.name, 'error': str(e)},
                'source': str(file_path)
            }

    def load_directory(
        self,
        directory: Union[str, Path],
        recursive: bool = True,
        file_patterns: Optional[List[str]] = None
    ) -> List[Dict[str, Any]]:
        """Load all supported documents from a directory"""
        directory = Path(directory)

        if not directory.exists():
            raise FileNotFoundError(f"Directory not found: {directory}")

        documents = []

        if recursive:
            pattern = "**/*"
        else:
            pattern = "*"

        for file_path in directory.glob(pattern):
            if file_path.is_file():
                # Check file patterns if specified
                if file_patterns:
                    if not any(file_path.match(pat) for pat in file_patterns):
                        continue

                # Check if extension is supported
                if file_path.suffix.lower() in self.supported_extensions:
                    doc = self.load_document(file_path)
                    documents.append(doc)

        console.print(f"[green]Loaded {len(documents)} documents from {directory}[/green]")
        return documents

    def _load_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        try:
            reader = PdfReader(file_path)
            text_parts = []

            for page in reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text)

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            return ""

    def _load_docx(self, file_path: Path) -> str:
        """Extract text from DOCX/DOC"""
        try:
            doc = Document(file_path)
            text_parts = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            return ""

    def _load_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX/PPT"""
        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_parts.append(" ".join(slide_text))

            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            return ""

    def _load_excel(self, file_path: Path) -> str:
        """Extract text from Excel files"""
        try:
            # Try with openpyxl first
            try:
                workbook = load_workbook(file_path, data_only=True)
                text_parts = []

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_text = [f"Sheet: {sheet_name}"]

                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            sheet_text.append(" | ".join(row_text))

                    text_parts.append("\n".join(sheet_text))

                return "\n\n".join(text_parts)

            except:
                # Fallback to pandas
                df = pd.read_excel(file_path, sheet_name=None)
                text_parts = []

                for sheet_name, data in df.items():
                    sheet_text = [f"Sheet: {sheet_name}"]
                    sheet_text.append(data.to_string())
                    text_parts.append("\n".join(sheet_text))

                return "\n\n".join(text_parts)

        except Exception as e:
            logger.error(f"Error reading Excel {file_path}: {e}")
            return ""

    def _load_csv(self, file_path: Path) -> str:
        """Extract text from CSV"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error reading CSV {file_path}: {e}")
            return ""

    def _load_text(self, file_path: Path) -> str:
        """Load plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Error reading text file {file_path}: {e}")
                return ""
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ""

    def _load_markdown(self, file_path: Path) -> str:
        """Load and convert Markdown to text"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()

            # Convert markdown to HTML then to text
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            logger.error(f"Error reading Markdown {file_path}: {e}")
            return ""

    def _load_html(self, file_path: Path) -> str:
        """Extract text from HTML"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()

            return soup.get_text()
        except Exception as e:
            logger.error(f"Error reading HTML {file_path}: {e}")
            return ""

    def _load_json(self, file_path: Path) -> str:
        """Extract text from JSON"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            # Convert JSON to readable text
            return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"Error reading JSON {file_path}: {e}")
            return ""

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions"""
        return list(self.supported_extensions.keys())

    def is_supported(self, file_path: Union[str, Path]) -> bool:
        """Check if file type is supported"""
        extension = Path(file_path).suffix.lower()
        return extension in self.supported_extensions